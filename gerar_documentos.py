import json
import os
import sys

import pandas as pd
from docx import Document
from docx.shared import Pt
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches


# ── xlsx ──────────────────────────────────────────────────────────────────────

def criar_xlsx(caminho: str, doc: dict) -> None:
    with pd.ExcelWriter(caminho, engine="openpyxl") as writer:
        if "entradas" in doc:
            # modo log: aba única tabular
            colunas = doc["colunas"]
            linhas = [[e.get(c, "") for c in colunas] for e in doc["entradas"]]
            df = pd.DataFrame(linhas, columns=colunas)
            df.to_excel(writer, sheet_name="Log", index=False)
        else:
            for aba in doc["abas"]:
                df = pd.DataFrame(aba["linhas"], columns=aba["colunas"])
                df.to_excel(writer, sheet_name=aba["nome"], index=False)


# ── docx ──────────────────────────────────────────────────────────────────────

def criar_docx(caminho: str, doc: dict) -> None:
    d = Document()
    for secao in doc.get("secoes", []):
        d.add_heading(secao["titulo"], level=1)
        d.add_paragraph(secao["conteudo"])

    campos = doc.get("campos_criticos", {})
    if campos:
        d.add_heading("Campos Críticos", level=2)
        tabela = d.add_table(rows=1, cols=2)
        tabela.style = "Table Grid"
        tabela.rows[0].cells[0].text = "Campo"
        tabela.rows[0].cells[1].text = "Valor"
        for chave, valor in campos.items():
            linha = tabela.add_row().cells
            linha[0].text = chave
            linha[1].text = str(valor)

    d.save(caminho)


# ── pdf ───────────────────────────────────────────────────────────────────────

_FONT_REGULAR = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
_FONT_BOLD    = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
_FONT_ITALIC  = "/usr/share/fonts/truetype/liberation/LiberationSans-Italic.ttf"


class _PDF(FPDF):
    def __init__(self):
        super().__init__()
        self.add_font("LiberationSans", fname=_FONT_REGULAR)
        self.add_font("LiberationSans", style="B", fname=_FONT_BOLD)
        self.add_font("LiberationSans", style="I", fname=_FONT_ITALIC)


def criar_pdf(caminho: str, doc: dict) -> None:
    pdf = _PDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("LiberationSans", size=11)

    if "entradas" in doc:
        # modo log: tabela em página única
        pdf.add_page()
        colunas = doc["colunas"]
        col_w = max(10, (pdf.w - 20) / max(len(colunas), 1))

        pdf.set_font("LiberationSans", style="B", size=10)
        for col in colunas:
            pdf.cell(col_w, 8, str(col)[:30], border=1)
        pdf.ln()

        pdf.set_font("LiberationSans", size=9)
        for entrada in doc["entradas"]:
            for col in colunas:
                pdf.cell(col_w, 7, str(entrada.get(col, ""))[:30], border=1)
            pdf.ln()
    else:
        secoes = doc.get("secoes", [])
        pdf.add_page()
        for secao in secoes:
            pdf.set_font("LiberationSans", style="B", size=12)
            pdf.multi_cell(0, 8, secao["titulo"], new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("LiberationSans", size=11)
            pdf.multi_cell(0, 7, secao["conteudo"], new_x="LMARGIN", new_y="NEXT")
            pdf.ln(3)

        assinatura = doc.get("assinatura", "")
        data_emissao = doc.get("data_emissao", "")
        if assinatura or data_emissao:
            pdf.ln(5)
            pdf.set_font("LiberationSans", style="I", size=10)
            if assinatura:
                pdf.cell(0, 6, f"Assinatura: {assinatura}", new_x="LMARGIN", new_y="NEXT")
            if data_emissao:
                pdf.cell(0, 6, f"Data de emissão: {data_emissao}", new_x="LMARGIN", new_y="NEXT")

    pdf.output(caminho)


# ── pptx ──────────────────────────────────────────────────────────────────────

def criar_pptx(caminho: str, doc: dict) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # "Title and Content"

    for slide_data in doc.get("slides", []):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("titulo", "")
        corpo = slide.placeholders[1]
        corpo.text = slide_data.get("conteudo", "")
        # dados descartados silenciosamente

    prs.save(caminho)


# ── eml ───────────────────────────────────────────────────────────────────────

def criar_eml(caminho: str, doc: dict) -> None:
    cab = doc.get("cabecalho", {})
    linhas = [
        f"De: {cab.get('de', '')}",
        f"Para: {cab.get('para', '')}",
        f"Assunto: {cab.get('assunto', '')}",
        f"Data: {cab.get('data', '')}",
        "",
        doc.get("corpo", ""),
    ]
    anexos = doc.get("anexos") or []
    if anexos:
        linhas += ["", f"Anexos: {', '.join(anexos)}"]

    with open(caminho, "w", encoding="utf-8") as f:
        f.write("\n".join(linhas))


# ── dispatcher ────────────────────────────────────────────────────────────────

HANDLERS = {
    "xlsx": criar_xlsx,
    "docx": criar_docx,
    "pdf":  criar_pdf,
    "pptx": criar_pptx,
    "eml":  criar_eml,
}


def processar_documentos(caminho_json: str) -> None:
    with open(caminho_json, "r", encoding="utf-8") as f:
        dados = json.load(f)  # exceção propagada se malformado

    diretorio_saida = os.path.join(
        os.path.dirname(os.path.abspath(caminho_json)), "documentos_gerados"
    )
    os.makedirs(diretorio_saida, exist_ok=True)

    for doc in dados["documentos"]:
        tipo = doc.get("tipo", "")
        handler = HANDLERS.get(tipo)
        if not handler:
            print(f"[AVISO] Tipo não suportado: '{tipo}' — pulando '{doc.get('arquivo', '?')}'")
            continue

        caminho = os.path.join(diretorio_saida, doc["arquivo"])
        try:
            handler(caminho, doc)
            print(f"Criado: {caminho}")
        except Exception as e:
            print(f"[ERRO] Falha ao gerar '{doc.get('arquivo', '?')}': {e}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python gerar_documentos.py <caminho/documentos.json>")
        sys.exit(1)

    processar_documentos(sys.argv[1])
